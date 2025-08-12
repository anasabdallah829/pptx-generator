import streamlit as st
import zipfile
import os
import io
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER, MSO_SHAPE_TYPE
import shutil
from pptx.util import Inches
import random

# إعداد صفحة Streamlit
st.set_page_config(page_title="PowerPoint Image Replacer", layout="centered")
st.title("🔄 PowerPoint Image & Placeholder Replacer")
st.markdown("---")

# واجهة المستخدم لرفع الملفات
uploaded_pptx = st.file_uploader("📂 اختر ملف PowerPoint (.pptx)", type=["pptx"], key="pptx_uploader")
uploaded_zip = st.file_uploader("🗜️ اختر ملف ZIP يحتوي على مجلدات صور", type=["zip"], key="zip_uploader")

# خيارات المعالجة
st.markdown("### ⚙️ إعدادات المعالجة")
image_order_option = st.radio(
    "كيف تريد ترتيب الصور في الشرائح؟",
    ("بالترتيب (افتراضي)", "عشوائي"),
    index=0
)

# إنشاء قائمة لحفظ التفاصيل
if 'processing_details' not in st.session_state:
    st.session_state.processing_details = []

def add_detail(message, detail_type="info"):
    """إضافة تفصيل جديد إلى قائمة التفاصيل"""
    st.session_state.processing_details.append({
        'message': message,
        'type': detail_type
    })

def clear_details():
    """مسح جميع التفاصيل"""
    st.session_state.processing_details = []

def show_details_section():
    """عرض قسم التفاصيل"""
    if st.session_state.processing_details:
        with st.expander("📋 تفاصيل المعالجة", expanded=False):
            for detail in st.session_state.processing_details:
                if detail['type'] == 'success':
                    st.success(detail['message'])
                elif detail['type'] == 'warning':
                    st.warning(detail['message'])
                elif detail['type'] == 'error':
                    st.error(detail['message'])
                else:
                    st.info(detail['message'])


def analyze_first_slide(prs):
    """
    تحليل الشريحة الأولى: إرجاع نتائج حتى لو لم توجد مواضع للصور.
    """
    if len(prs.slides) == 0:
        return False, "لا توجد شرائح في الملف"

    first_slide = prs.slides[0]
    
    picture_placeholders = [
        shape for shape in first_slide.shapes
        if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.PICTURE
    ]
    
    # استخدام نفس طريقة الكود المرجعي للصور العادية
    regular_pictures = [
        shape for shape in first_slide.shapes 
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    ]
    
    total_image_slots = len(picture_placeholders) + len(regular_pictures)

    return True, {
        'placeholders': len(picture_placeholders),
        'regular_pictures': len(regular_pictures),
        'total_slots': total_image_slots,
        'slide_layout': first_slide.slide_layout
    }


def get_shape_formatting(shape):
    """
    استخراج تنسيقات الشكل الأصلية
    """
    formatting = {
        'left': shape.left,
        'top': shape.top,
        'width': shape.width,
        'height': shape.height,
        'rotation': getattr(shape, 'rotation', 0),
    }
    
    # استخراج تنسيقات إضافية إذا كانت متوفرة
    try:
        if hasattr(shape, 'shadow'):
            formatting['shadow'] = {
                'inherit': shape.shadow.inherit,
                'visible': getattr(shape.shadow, 'visible', None)
            }
    except:
        pass
    
    try:
        if hasattr(shape, 'line'):
            formatting['line'] = {
                'color': getattr(shape.line.color, 'rgb', None),
                'width': getattr(shape.line, 'width', None)
            }
    except:
        pass
    
    try:
        if hasattr(shape, 'fill'):
            formatting['fill_type'] = getattr(shape.fill, 'type', None)
    except:
        pass
    
    return formatting


def apply_shape_formatting(new_shape, formatting):
    """
    تطبيق التنسيقات على الشكل الجديد
    """
    try:
        # تطبيق الموقع والحجم
        new_shape.left = formatting['left']
        new_shape.top = formatting['top']
        new_shape.width = formatting['width']
        new_shape.height = formatting['height']
        
        # تطبيق الدوران
        if 'rotation' in formatting and formatting['rotation'] != 0:
            new_shape.rotation = formatting['rotation']
        
        # تطبيق الظل
        if 'shadow' in formatting:
            try:
                if formatting['shadow']['visible'] is not None:
                    new_shape.shadow.visible = formatting['shadow']['visible']
            except:
                pass
        
        # تطبيق الحدود
        if 'line' in formatting:
            try:
                if formatting['line']['width'] is not None:
                    new_shape.line.width = formatting['line']['width']
                if formatting['line']['color'] is not None:
                    new_shape.line.color.rgb = formatting['line']['color']
            except:
                pass
                
    except Exception as e:
        # في حالة فشل تطبيق أي تنسيق، نتجاهل الخطأ ونكمل
        pass


def get_image_shapes_info(slide):
    """
    استخراج معلومات مفصلة عن أشكال الصور من الشريحة مع التنسيقات
    """
    image_shapes_info = []
    
    # البحث عن placeholders للصور
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
            formatting = get_shape_formatting(shape)
            image_shapes_info.append({
                'shape': shape,
                'type': 'placeholder',
                'formatting': formatting,
                'position': (shape.top, shape.left)
            })
    
    # البحث عن الصور العادية باستخدام نفس طريقة الكود المرجعي
    regular_pictures = [
        shape for shape in slide.shapes 
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    ]
    
    # إضافة معلومات الصور العادية مع التنسيقات
    for shape in regular_pictures:
        formatting = get_shape_formatting(shape)
        image_shapes_info.append({
            'shape': shape,
            'type': 'picture',
            'formatting': formatting,
            'position': (shape.top, shape.left)
        })
    
    # ترتيب حسب الموقع (من الأعلى للأسفل، من اليسار لليمين)
    image_shapes_info.sort(key=lambda x: x['position'])
    return image_shapes_info


def get_template_image_positions(slide):
    """
    استخراج مواقع الصور من القالب مع التنسيقات الكاملة
    """
    image_positions = []
    
    # استخدام نفس الطريقة من الكود المرجعي مع إضافة التنسيقات
    image_shapes = [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
    for shape in image_shapes:
        formatting = get_shape_formatting(shape)
        image_positions.append(formatting)
    
    # إضافة placeholders أيضاً
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
            formatting = get_shape_formatting(shape)
            image_positions.append(formatting)
    
    return image_positions


def replace_image_in_shape(slide, shape_info, image_path):
    """
    استبدال صورة في شكل محدد مع الحفاظ على التنسيقات الأصلية
    """
    try:
        shape = shape_info['shape']
        shape_type = shape_info['type']
        original_formatting = shape_info['formatting']
        
        if shape_type == 'placeholder':
            # معالجة placeholders بالطريقة العادية
            try:
                with open(image_path, 'rb') as img_file:
                    shape.insert_picture(img_file)
                add_detail(f"✅ تم استبدال placeholder بنجاح: {os.path.basename(image_path)}", "success")
                return True
            except Exception as e:
                add_detail(f"⚠ فشل في استبدال placeholder، محاولة طريقة بديلة: {e}", "warning")
                
                # طريقة بديلة للـ placeholders مع الحفاظ على التنسيقات
                try:
                    # حذف الشكل القديم
                    shape_element = shape._element
                    shape_element.getparent().remove(shape_element)
                    
                    # إضافة صورة جديدة مع التنسيقات
                    new_shape = slide.shapes.add_picture(
                        image_path, 
                        original_formatting['left'], 
                        original_formatting['top'], 
                        original_formatting['width'], 
                        original_formatting['height']
                    )
                    
                    # تطبيق التنسيقات الأصلية
                    apply_shape_formatting(new_shape, original_formatting)
                    
                    add_detail(f"✅ تم استبدال placeholder بالطريقة البديلة مع الحفاظ على التنسيقات: {os.path.basename(image_path)}", "success")
                    return True
                except Exception as e2:
                    add_detail(f"❌ فشل في استبدال placeholder: {e2}", "error")
                    return False
        
        elif shape_type == 'picture':
            # استبدال الصور العادية مع الحفاظ على التنسيقات
            try:
                # حذف الصورة القديمة
                shape_element = shape._element
                shape_element.getparent().remove(shape_element)
                
                # إضافة الصورة الجديدة مع التنسيقات الأصلية
                new_shape = slide.shapes.add_picture(
                    image_path, 
                    original_formatting['left'], 
                    original_formatting['top'], 
                    original_formatting['width'], 
                    original_formatting['height']
                )
                
                # تطبيق التنسيقات الأصلية
                apply_shape_formatting(new_shape, original_formatting)
                
                add_detail(f"✅ تم استبدال الصورة العادية مع الحفاظ على التنسيقات: {os.path.basename(image_path)}", "success")
                return True
            except Exception as e:
                add_detail(f"❌ فشل في استبدال الصورة العادية: {e}", "error")
                return False
        
        return False
        
    except Exception as e:
        add_detail(f"❌ خطأ عام في استبدال الصورة: {e}", "error")
        return False


def add_images_using_template_positions(slide, images, image_positions):
    """
    إضافة الصور باستخدام مواقع القالب مع الحفاظ على التنسيقات
    """
    added_count = 0
    
    for idx, formatting in enumerate(image_positions):
        if idx < len(images):
            try:
                # إضافة الصورة مع التنسيقات الأصلية
                new_shape = slide.shapes.add_picture(
                    images[idx], 
                    formatting['left'], 
                    formatting['top'], 
                    formatting['width'], 
                    formatting['height']
                )
                
                # تطبيق التنسيقات الأصلية
                apply_shape_formatting(new_shape, formatting)
                
                added_count += 1
                add_detail(f"✅ تم إضافة صورة بطريقة القالب مع التنسيقات: {os.path.basename(images[idx])}", "success")
            except Exception as e:
                add_detail(f"❌ فشل في إضافة صورة: {e}", "error")
    
    return added_count


def add_title_to_slide(slide, folder_name):
    """
    إضافة أو تحديث عنوان الشريحة
    """
    try:
        # البحث عن placeholder للعنوان
        title_shapes = [
            shape for shape in slide.shapes
            if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.TITLE
        ]
        
        if title_shapes:
            # تحديث العنوان الموجود
            title_shapes[0].text = folder_name
            add_detail(f"✅ تم تحديث العنوان: {folder_name}", "success")
        else:
            # إضافة عنوان جديد
            try:
                textbox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
                text_frame = textbox.text_frame
                text_frame.text = folder_name
                
                # تنسيق النص
                paragraph = text_frame.paragraphs[0]
                paragraph.font.size = Inches(0.4)
                paragraph.font.bold = True
                
                add_detail(f"✅ تم إضافة عنوان جديد: {folder_name}", "success")
            except Exception as e:
                add_detail(f"⚠ فشل في إضافة العنوان: {e}", "warning")
    except Exception as e:
        add_detail(f"⚠ خطأ في معالجة العنوان: {e}", "warning")


def process_folder_images(slide, folder_path, folder_name, template_shapes_info, template_positions, mismatch_action):
    """
    معالجة صور مجلد واحد وإضافتها للشريحة مع الحفاظ على التنسيقات
    """
    # الحصول على قائمة الصور
    imgs = [f for f in os.listdir(folder_path) 
            if f.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.webp'))]
    
    if not imgs:
        add_detail(f"⚠ المجلد {folder_name} فارغ من الصور", "warning")
        return 0
    
    # ترتيب الصور بناءً على اختيار المستخدم
    if image_order_option == "عشوائي":
        random.shuffle(imgs)
        add_detail(f"🔀 تم ترتيب صور المجلد {folder_name} عشوائياً", "info")
    else:
        imgs.sort()
        add_detail(f"📋 تم ترتيب صور المجلد {folder_name} أبجدياً", "info")
    
    # تحويل أسماء الصور إلى مسارات كاملة
    image_paths = [os.path.join(folder_path, img) for img in imgs]
    
    # إضافة العنوان
    add_title_to_slide(slide, folder_name)
    
    # الحصول على معلومات أشكال الصور في الشريحة الجديدة
    new_shapes_info = get_image_shapes_info(slide)
    
    replaced_count = 0
    
    if new_shapes_info:
        # إذا وجدت أشكال صور في الشريحة الجديدة، استبدلها
        add_detail(f"📸 وجدت {len(new_shapes_info)} شكل صورة في الشريحة الجديدة", "info")
        
        # معالجة اختلاف عدد الصور
        if mismatch_action == 'skip_folder' and len(imgs) != len(new_shapes_info):
            add_detail(f"ℹ تم تخطي المجلد {folder_name} لوجود اختلاف في عدد الصور", "info")
            return 0
        
        # استبدال الصور
        for i, shape_info in enumerate(new_shapes_info):
            if mismatch_action == 'truncate' and i >= len(imgs):
                break
            
            # اختيار الصورة (مع التكرار إذا لزم الأمر)
            image_path = image_paths[i % len(image_paths)]
            
            # التحقق من وجود الملف
            if not os.path.exists(image_path):
                add_detail(f"⚠ الملف غير موجود: {image_path}", "warning")
                continue
            
            # استبدال الصورة مع الحفاظ على التنسيقات
            success = replace_image_in_shape(slide, shape_info, image_path)
            if success:
                replaced_count += 1
    
    elif template_positions:
        # إذا لم توجد أشكال في الشريحة الجديدة، استخدم مواقع القالب
        add_detail(f"📍 استخدام مواقع القالب ({len(template_positions)} موقع)", "info")
        
        replaced_count = add_images_using_template_positions(
            slide, image_paths, template_positions
        )
    
    else:
        # إضافة الصورة الأولى في موقع افتراضي
        add_detail(f"⚠ لا توجد مواضع للصور، إضافة الصورة الأولى في موقع افتراضي", "warning")
        
        if image_paths:
            try:
                slide.shapes.add_picture(image_paths[0], Inches(1), Inches(2), Inches(8), Inches(5))
                add_detail(f"✅ تم إضافة الصورة الأولى في موقع افتراضي: {imgs[0]}", "success")
                replaced_count = 1
            except Exception as e:
                add_detail(f"❌ فشل في إضافة الصورة الافتراضية: {e}", "error")
    
    return replaced_count


def main():
    if uploaded_pptx and uploaded_zip:
        if "process_started" not in st.session_state:
            st.session_state.process_started = False

        if st.button("🚀 بدء المعالجة") or st.session_state.process_started:
            st.session_state.process_started = True
            
            # مسح التفاصيل السابقة
            clear_details()
            
            temp_dir = None
            try:
                st.info("📦 جاري فحص الملفات...")
                add_detail("📦 بدء فحص الملفات المرفوعة", "info")
                
                zip_bytes = io.BytesIO(uploaded_zip.read())
                with zipfile.ZipFile(zip_bytes, "r") as zip_ref:
                    temp_dir = "temp_images"
                    if os.path.exists(temp_dir):
                        shutil.rmtree(temp_dir)
                    os.makedirs(temp_dir)
                    zip_ref.extractall(temp_dir)
                
                add_detail("📂 تم استخراج الملف المضغوط بنجاح", "success")
                
                all_items = os.listdir(temp_dir)
                folder_paths = []
                
                for item in all_items:
                    item_path = os.path.join(temp_dir, item)
                    if os.path.isdir(item_path):
                        # التحقق من وجود صور في المجلد
                        imgs_in_folder = [f for f in os.listdir(item_path) 
                                        if f.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.webp'))]
                        if imgs_in_folder:
                            folder_paths.append(item_path)
                            add_detail(f"📁 المجلد '{item}' يحتوي على {len(imgs_in_folder)} صورة", "info")
                
                if not folder_paths:
                    st.error("❌ لا توجد مجلدات تحتوي على صور في الملف المضغوط.")
                    add_detail("❌ لا توجد مجلدات تحتوي على صور في الملف المضغوط", "error")
                    show_details_section()
                    st.stop()
                
                folder_paths.sort()
                st.success(f"✅ تم العثور على {len(folder_paths)} مجلد يحتوي على صور")
                add_detail(f"✅ تم العثور على {len(folder_paths)} مجلد يحتوي على صور", "success")

                prs = Presentation(io.BytesIO(uploaded_pptx.read()))
                
                st.info("🔍 جاري تحليل الشريحة الأولى...")
                add_detail("🔍 بدء تحليل الشريحة الأولى", "info")
                
                ok, analysis_result = analyze_first_slide(prs)
                if not ok:
                    st.error(f"❌ {analysis_result}")
                    add_detail(f"❌ {analysis_result}", "error")
                    show_details_section()
                    st.stop()
                
                st.success("✅ تحليل الشريحة الأولى جاهز")
                add_detail("✅ تم تحليل الشريحة الأولى بنجاح", "success")
                
                col1, col2, col3 = st.columns(3)
                with col1: st.metric("Placeholders للصور", analysis_result['placeholders'])
                with col2: st.metric("الصور العادية", analysis_result['regular_pictures'])
                with col3: st.metric("إجمالي أماكن الصور", analysis_result['total_slots'])
                
                add_detail(f"📊 تفاصيل التحليل: {analysis_result['placeholders']} placeholders، {analysis_result['regular_pictures']} صور عادية، {analysis_result['total_slots']} إجمالي", "info")
                
                first_slide = prs.slides[0]
                template_shapes_info = get_image_shapes_info(first_slide)
                template_positions = get_template_image_positions(first_slide)
                
                if not template_shapes_info and not template_positions:
                    st.warning("⚠ الشريحة الأولى لا تحتوي على مواضع صور. سيتم إضافة الصورة الأولى من كل مجلد فقط.")
                    add_detail("⚠ الشريحة الأولى لا تحتوي على مواضع صور", "warning")
                    slide_layout = prs.slide_layouts[6]  # Blank layout
                else:
                    slide_layout = analysis_result['slide_layout']

                # فحص التطابق في عدد الصور
                expected_count = max(len(template_shapes_info), len(template_positions))
                mismatch_folders = []
                for fp in folder_paths:
                    imgs = [f for f in os.listdir(fp) 
                           if f.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.webp'))]
                    if len(imgs) != expected_count:
                        mismatch_folders.append((os.path.basename(fp), len(imgs), expected_count))
                
                if mismatch_folders and 'mismatch_action' not in st.session_state:
                    with st.form("mismatch_form"):
                        st.warning("⚠ تم اكتشاف اختلاف في عدد الصور لبعض المجلدات مقارنة بعدد مواضع الصور في الشريحة الأولى.")
                        for name, img_count, _ in mismatch_folders:
                            st.write(f"- المجلد `{name}` يحتوي على {img_count} صورة.")
                            add_detail(f"⚠ المجلد '{name}' يحتوي على {img_count} صورة بدلاً من {expected_count}", "warning")
                        st.markdown(f"**عدد مواضع الصور في القالب: {expected_count}**")

                        choice_text = st.radio(
                            "اختر كيف تريد التعامل مع المجلدات التي يختلف عدد صورها:",
                            ("استبدال فقط حتى أقل عدد (truncate)", "تكرار الصور لملء جميع المواضع (repeat)", "تخطي المجلدات ذات الاختلاف (skip_folder)", "إيقاف العملية (stop)"),
                            index=0
                        )
                        submit_choice = st.form_submit_button("✅ تأكيد الاختيار والمتابعة")

                    if submit_choice:
                        if choice_text.startswith("استبدال فقط"): 
                            st.session_state['mismatch_action'] = 'truncate'
                            add_detail("⚙️ تم اختيار: استبدال فقط حتى أقل عدد", "info")
                        elif choice_text.startswith("تكرار"): 
                            st.session_state['mismatch_action'] = 'repeat'
                            add_detail("⚙️ تم اختيار: تكرار الصور لملء جميع المواضع", "info")
                        elif choice_text.startswith("تخطي"): 
                            st.session_state['mismatch_action'] = 'skip_folder'
                            add_detail("⚙️ تم اختيار: تخطي المجلدات ذات الاختلاف", "info")
                        else: 
                            st.session_state['mismatch_action'] = 'stop'
                            add_detail("⚙️ تم اختيار: إيقاف العملية", "info")
                    else:
                        show_details_section()
                        st.stop()
                
                if 'mismatch_action' in st.session_state:
                    mismatch_action = st.session_state['mismatch_action']
                else:
                    mismatch_action = 'truncate'

                if mismatch_action == 'stop':
                    st.error("❌ تم إيقاف العملية بناءً على اختيار المستخدم.")
                    add_detail("❌ تم إيقاف العملية بناءً على اختيار المستخدم", "error")
                    show_details_section()
                    st.stop()

                st.info("🔄 جاري إضافة الشرائح الجديدة...")
                add_detail("🔄 بدء إضافة الشرائح الجديدة", "info")
                
                total_replaced = 0
                created_slides = 0

                progress_bar = st.progress(0)
                status_text = st.empty()

                for folder_idx, folder_path in enumerate(folder_paths):
                    folder_name = os.path.basename(folder_path)
                    status_text.text(f"🔄 معالجة المجلد {folder_idx + 1}/{len(folder_paths)}: {folder_name}")
                    add_detail(f"🔄 بدء معالجة المجلد: {folder_name}", "info")

                    try:
                        # إنشاء شريحة جديدة
                        new_slide = prs.slides.add_slide(slide_layout)
                        created_slides += 1
                        add_detail(f"📄 تم إنشاء شريحة جديدة للمجلد: {folder_name}", "success")
                        
                        # معالجة صور المجلد مع الحفاظ على التنسيقات
                        replaced_count = process_folder_images(
                            new_slide, folder_path, folder_name, 
                            template_shapes_info, template_positions, mismatch_action
                        )
                        
                        total_replaced += replaced_count
                        add_detail(f"✅ تم إنشاء شريحة للمجلد '{folder_name}' واستبدال {replaced_count} صورة", "success")
                    
                    except Exception as e:
                        st.error(f"❌ خطأ في معالجة المجلد {folder_name}: {e}")
                        add_detail(f"❌ خطأ في معالجة المجلد {folder_name}: {e}", "error")

                    progress_bar.progress((folder_idx + 1) / len(folder_paths))

                progress_bar.empty()
                status_text.empty()

                st.success("🎉 تم الانتهاء من المعالجة مع الحفاظ على جميع التنسيقات!")
                add_detail("🎉 تم الانتهاء من المعالجة بنجاح", "success")
                
                # تنظيف session state
                if 'mismatch_action' in st.session_state: 
                    del st.session_state['mismatch_action']
                if 'process_started' in st.session_state: 
                    del st.session_state['process_started']

                col1, col2, col3 = st.columns(3)
                with col1: st.metric("الشرائح المُضافة", created_slides)
                with col2: st.metric("الصور المُستبدلة", total_replaced)
                with col3: st.metric("المجلدات المُعالجة", len(folder_paths))

                add_detail(f"📈 النتائج النهائية: {created_slides} شريحة، {total_replaced} صورة مُستبدلة، {len(folder_paths)} مجلد مُعالج", "success")

                if created_slides == 0:
                    st.error("❌ لم يتم إضافة أي شرائح.")
                    add_detail("❌ لم يتم إضافة أي شرائح", "error")
                    show_details_section()
                    st.stop()

                # حفظ الملف
                original_name = os.path.splitext(uploaded_pptx.name)[0]
                output_filename = f"{original_name}_Updated.pptx"
                output_buffer = io.BytesIO()
                prs.save(output_buffer)
                output_buffer.seek(0)

                st.success(f"✅ تم إنشاء ملف PowerPoint جديد بـ {created_slides} شريحة مع الحفاظ على جميع التنسيقات!")
                add_detail(f"💾 تم حفظ الملف: {output_filename}", "success")

                st.download_button(
                    label="⬇️ تحميل الملف المُحدث",
                    data=output_buffer.getvalue(),
                    file_name=output_filename,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    key="download_button"
                )

                # عرض قسم التفاصيل بعد اكتمال المعالجة
                show_details_section()

            except Exception as e:
                st.error(f"❌ خطأ أثناء المعالجة: {e}")
                add_detail(f"❌ خطأ عام أثناء المعالجة: {e}", "error")
                show_details_section()
            finally:
                if temp_dir and os.path.exists(temp_dir):
                    try:
                        shutil.rmtree(temp_dir)
                        add_detail("🧹 تم تنظيف الملفات المؤقتة", "info")
                    except Exception as cleanup_error:
                        add_detail(f"⚠ خطأ في تنظيف الملفات المؤقتة: {cleanup_error}", "warning")
    else:
        st.info("📋 يُرجى رفع ملف PowerPoint وملف ZIP للبدء")

        with st.expander("📖 تعليمات الاستخدام"):
            st.markdown("""
            ### كيفية استخدام التطبيق:

            1.  **ملف PowerPoint (.pptx):**
                - يجب أن يحتوي على شريحة واحدة على الأقل.
                - يتم استخدام تنسيق الشريحة الأولى كقالب.
                - **يتم الحفاظ على جميع التنسيقات الأصلية للصور** (الحجم، الموقع، الدوران، الظلال، الحدود، إلخ).

            2.  **ملف ZIP:**
                - يجب أن يحتوي على مجلدات، وكل مجلد يحتوي على صور.
                - أسماء المجلدات ستصبح عناوين الشرائح.

            3.  **النتيجة:**
                - شريحة منفصلة لكل مجلد.
                - يتم استبدال الصور و placeholders في القالب بصور من المجلدات.
                - في حال عدم وجود مواضع للصور في القالب، تُضاف الصورة الأولى من كل مجلد.
                - **جميع التنسيقات الأصلية محفوظة**.

            ### أنواع الصور المدعومة:
            - PNG, JPG, JPEG, GIF, BMP, TIFF, WEBP

            ### الميزات الجديدة:
            - 🎨 **الحفاظ على التنسيقات**: جميع تنسيقات الصور الأصلية (الحجم، الموقع، الدوران، الظلال، الحدود) محفوظة
            - 📋 **تفاصيل المعالجة**: يمكن عرض تفاصيل كاملة لعملية المعالجة
            - 🔀 **ترتيب الصور**: اختيار بين الترتيب الأبجدي أو العشوائي
            - ⚙️ **خيارات مرونة**: التعامل مع اختلاف عدد الصور بطرق متعددة
            """)

if __name__ == '__main__':
    main()
